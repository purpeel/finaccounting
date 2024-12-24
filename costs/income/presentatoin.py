from pptx import Presentation

# Создаем презентацию
presentation = Presentation()

# Список содержания слайдов
slides_content = [
    {"title": "Концептуализация цивилизации",
     "content": "Взгляд Сэмюэля Хантингтона\n\nПрезентация подготовлена: [Ваше имя]\nДата: [Дата]"},

    {"title": "Введение",
     "content": "• Цивилизация: культурно-историческая общность людей\n"
                "• Почему важно изучать цивилизации?\n"
                "• Книга Хантингтона 'Столкновение цивилизаций' и ее влияние"},

    {"title": "Что такое цивилизация?",
     "content": "• Определение: большая культурная группа с общей историей, традициями и ценностями\n"
                "• Характеристики:\n"
                "  - Язык\n"
                "  - Религия\n"
                "  - Культура\n"
                "  - Историческое наследие"},

    {"title": "Исторические примеры цивилизаций",
     "content": "• Египетская, Месопотамская, Греческая, Китайская\n"
                "• Современные цивилизации:\n"
                "  - Западная\n"
                "  - Исламская\n"
                "  - Индийская\n"
                "  - Китайская"},

    {"title": "Цивилизационный подход Хантингтона",
     "content": "• Основные тезисы:\n"
                "  - 'Столкновение цивилизаций' — новая парадигма мирового порядка\n"
                "  - Культура и идентичность — основа цивилизаций\n"
                "• Цивилизационные блоки:\n"
                "  - Западная, Синская, Православная, Исламская и другие"},

    {"title": "Особенности и границы цивилизаций",
     "content": "• Объединяющие факторы:\n"
                "  - Общая религия, язык, культура\n"
                "• Разделяющие факторы:\n"
                "  - Конфликты, различия в ценностях и традициях"},

    {"title": "Плюсы цивилизационного подхода",
     "content": "• Учет культурной идентичности\n"
                "• Объяснение глобальных конфликтов\n"
                "• Анализ многополярного мира"},

    {"title": "Минусы цивилизационного подхода",
     "content": "• Упрощение и обобщение культур\n"
                "• Игнорирование внутренних различий\n"
                "• Усиление стереотипов и конфликтов"},

    {"title": "Современное значение теории",
     "content": "• Использование в геополитике\n"
                "• Альтернативы: глобализация, модернизация\n"
                "• Критика теории Хантингтона"},

    {"title": "Примеры столкновений цивилизаций",
     "content": "• Холодная война: идеологический и цивилизационный конфликт\n"
                "• Запад и исламский мир после 9/11"},

    {"title": "Заключение",
     "content": "• Основные выводы:\n"
                "  - Цивилизации — ключ к пониманию мировой политики\n"
                "• Теория Хантингтона остается актуальной\n"
                "• Вопросы для обсуждения:"},

    {"title": "Вопросы и обсуждение",
     "content": "• Как вы оцениваете идею 'столкновения цивилизаций'?\n"
                "• Какие аспекты теории Хантингтона вызывают сомнения?"}
]

# Создание слайдов
for slide in slides_content:
    slide_layout = presentation.slide_layouts[1]  # Заголовок и текст
    slide_obj = presentation.slides.add_slide(slide_layout)
    slide_obj.shapes.title.text = slide["title"]
    slide_obj.placeholders[1].text = slide["content"]

# Сохранение презентации
presentation_path = "C:\Концептуализация_цивилизации_Хантингтон.pptx"
presentation.save(presentation_path)
presentation_path
